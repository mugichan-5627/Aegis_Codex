import pptx

def set_shape_text(shape, text):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    if not tf.paragraphs:
        return
    # Use paragraph 0, run 0 to preserve formatting
    p = tf.paragraphs[0]
    if p.runs:
        p.runs[0].text = text
        for run in p.runs[1:]:
            run.text = ""
    else:
        p.text = text

# Load presentation
prs = pptx.Presentation(r'c:\Users\Moosa\Downloads\Project_Doomsday.pptx')

# Mapping: (slide_idx, shape_idx) -> new_text
replacements = {
    # ================ SLIDE 1 (Title) ================
    (0, 1): "[ GOOGLE CLOUD RAPID AGENT HACKATHON ]",
    (0, 2): "SWARM ACTIVE",
    (0, 4): "SWARM ACTIVE",
    (0, 6): "THREAT: HIGH",
    (0, 13): "[!]",
    (0, 14): "DOOMSDAY RAPID AGENT",
    (0, 16): "SYSTEMIC BLACK SWAN SIMULATION COCKPIT",
    (0, 18): '"Every enterprise has a breaking point. We expose it in 15 seconds."',
    (0, 20): "4",
    (0, 21): "FRACTURE AGENTS",
    (0, 23): "8",
    (0, 24): "MCP INDEX VECTORS",
    (0, 26): "<15s",
    (0, 27): "EXECUTION TIME",
    (0, 29): "BUILT WITH:  GOOGLE GEMINI 3  |  ELASTIC MCP  |  ARIZE PHOENIX  |  STREAMLIT",
    (0, 30): "MOOSA  //  IIM KOZHIKODE  //  BATCH 29",

    # ================ SLIDE 2 (The Problem) ================
    (1, 0): "// 01  THE PROBLEM",
    (1, 1): "BLIND SPOTS DEVALUE ENTERPRISES",
    (1, 5): "[01] THE SPEED DECAY TRAP",
    (1, 6): "Traditional risk reviews take weeks. Under sudden geopolitical fracture, valuations compress in hours. Slow analysis is an invitation to capital destruction.",
    (1, 9): "[02] ELUSIVE CONTEXT CORRELATIONS",
    (1, 10): "Threat vectors are modeled in isolation, ignoring ground-truth macro indices. Geopolitical shocks cascading into basis point margin cuts remain completely unmodeled.",
    (1, 13): "[03] TELEMETRY BLACK BOXES",
    (1, 14): "Most agentic architectures are unmonitored loops. Without real-time telemetry tracing and token cost monitoring, public deployments run blind.",
    (1, 17): "THE GAP",
    (1, 18): "No system exists that identifies current geopolitical shocks — validates them through an adversarial swarm — grounds them with Elastic MCP — traces telemetry via Arize Phoenix — and models distressed WACC waterfalls. Instantly.",
    (1, 20): "90%+",
    (1, 21): "OF AGENT SYSTEMS LACK ADVOCATE-LEVEL ADVERSARIAL VALIDATION",
    (1, 22): "LEADING TO MASSIVE BIAS AND ECHO-CHAMBER VERDICTS",
    (1, 23): "— AI ALIGNMENT RESEARCH 2025",
    (1, 25): "DOOMSDAY RAPID AGENT — HACKATHON CONFIDENTIAL",

    # ================ SLIDE 3 (Architecture) ================
    (2, 0): "// 02  SYSTEM ARCHITECTURE",
    (2, 1): "HOW IT WORKS: 7-STEP RAPID PIPELINE",
    (2, 5): "01",
    (2, 6): "USER INPUT",
    (2, 7): "Ticker + Custom Chaos Index (\u03c7)",
    (2, 11): "02",
    (2, 12): "TELEMETRY HOOK",
    (2, 13): "Phoenix collector setup (Graceful fallback to local dashboard)",
    (2, 17): "03",
    (2, 18): "GROUNDED INGEST",
    (2, 19): "Baseline financials + multi-factor WACC loading",
    (2, 23): "04",
    (2, 24): "ELASTIC MCP QUERY",
    (2, 25): "Elastic MCP queries live macro indices + unstructured risk feeds",
    (2, 28): "05",
    (2, 29): "FRACTURE DEBATE",
    (2, 30): "Swarm Geopolitical vs Supply Chain debate, judged by Gemini-3",
    (2, 33): "06",
    (2, 34): "GEOSPATIAL COCKPIT",
    (2, 35): "Plotly interactive map showing shock nexuses & HQ spokes",
    (2, 38): "07",
    (2, 39): "VALUATION WATERFALL",
    (2, 40): "Beta-adjusted WACC degradation + direct basis point haircut cascade",
    (2, 42): "\u25c0  PHASE 1: MCP DATA GROUNDING",
    (2, 43): "\u25c0  PHASE 2: DEBATE & VALUATION CASCADE",
    (2, 45): "\u26a1  4 AI AGENTS IN SWARM  ·  GEMINI 3 VIA GOOGLE STUDIO  ·  ELASTIC MCP  ·  ARIZE PHOENIX TRACING",

    # ================ SLIDE 4 (Multi-Agent Swarm) ================
    (3, 0): "// 03  MULTI-AGENT SYSTEM",
    (3, 1): "4 SPECIALIZED AGENTS IN ADVERSARIAL SWARM",
    (3, 5): "GEOPOLITICAL ANALYST",
    (3, 6): "MACRO THREATS",
    (3, 8): "T=0.3",
    (3, 10): "Extracts raw geopolitical shock indicators, locates geospatial coordinates, and estimates macro-risk severity.",
    (3, 13): "SUPPLY CHAIN ANALYST",
    (3, 14): "MICRO CASCADE",
    (3, 16): "T=0.4",
    (3, 18): "Analyzes logistical bottlenecks, route fractures, and basis point impact on operational margins. Pushes severity up.",
    (3, 21): "FINANCIAL ANALYST",
    (3, 22): "VALUATION IMPACT",
    (3, 24): "T=0.4",
    (3, 26): "Argues market pricing models, calculates mitigating assets, and evaluates baseline balance sheet resiliency. Pushes severity down.",
    (3, 29): "BLACK SWAN JUDGE",
    (3, 30): "FINAL CONSENSUS",
    (3, 32): "T=0.2",
    (3, 34): "Synthesizes the adversarial dialectic tribunal. Computes final severity scores & maps to distressed valuation cascades.",
    (3, 37): "ARIZE TELEMETRY CLIENT",
    (3, 38): "OBSERVABILITY",
    (3, 40): "T=0.1",
    (3, 42): "Instruments trace spans, logs coordinates and parameters directly to Arize Phoenix collector or in-app console.",
    (3, 44): "\u25b8  Dialectic Framework: Swarm isolates bias through dedicated bull/bear geopolitical and supply chain debate",
    (3, 45): "\u25b8  Full Traceability: Phoenix logs every prompt token, node latency, and agent reasoning pathway in real-time",

    # ================ SLIDE 5 (Differentiators) ================
    (4, 0): "// 04  DIFFERENTIATORS",
    (4, 1): "WHAT MAKES THIS DIFFERENT",
    (4, 5): "D1",
    (4, 6): "DYNAMIC MATH SHIFT",
    (4, 7): "Three distinct distress engines. Automatically shifts between Multi-Factor DCF, EV/Revenue, and Cyclical models. Computes distressed WACC increases.",
    (4, 10): "D2",
    (4, 11): "ELASTIC MCP GROUNDING",
    (4, 12): "Queries Elasticsearch index via MCP to validate macro feeds and historical precedents. Anchors LLM intelligence in verifiable facts.",
    (4, 15): "D3",
    (4, 16): "LIVE TRACE TELEMETRY",
    (4, 17): "Full telemetry visibility via Arize Phoenix. Graceful degradation routes spans to in-app dashboard card when external collectors are absent.",
    (4, 20): "D4",
    (4, 21): "MILITARY-GRADE COCKPIT",
    (4, 22): "Sleek high-contrast dark interface featuring responsive spoke-and-hub Plotly map connecting global shocks to company HQs.",
    (4, 25): "RAPID VS. ORIGINAL",
    (4, 27): "RAPID COCKPIT",
    (4, 28): "PROJECT DOOMSDAY",
    (4, 31): "Dynamic Swarm Debate",
    (4, 33): "\u2717",
    (4, 34): "Elastic MCP Grounding",
    (4, 36): "\u2717 (Standard Web Search Only)",
    (4, 38): "Arize Phoenix Tracing",
    (4, 40): "\u2717",
    (4, 41): "Telemetry Fallback",
    (4, 43): "\u2717",
    (4, 45): "Valuation Math Models",
    (4, 46): "3 (DCF, EV/Rev, Cyclical)",
    (4, 47): "5",
    (4, 48): "Analysis Latency",
    (4, 49): "<15 seconds",
    (4, 50): "<60 seconds",
    (4, 52): "HQ Spoke-and-Hub Map",
    (4, 54): "\u2713",
    (4, 55): "Telemetry Config",
    (4, 56): "Sidebar UI & Env",
    (4, 57): "Env variables only",

    # ================ SLIDE 6 (Go-To-Market) ================
    (5, 0): "// 05  GO-TO-MARKET",
    (5, 1): "FROM HACKATHON TO PRODUCT",
    (5, 32): "Custom domains  ·  Telemetry Sync  ·  White-label",
    (5, 42): "GEMINI API COST / ANALYSIS   $0.02\u20130.05",
    (5, 43): "ELASTIC / TAVILY SCANS     $0.01\u20130.02",
    (5, 45): "ARIZE TELEMETRY SYNC \u2014 FIRST-CLASS OBSERVABILITY AND LATENCY MONITORING \u2014 ENTERPRISE GRADE",

    # ================ SLIDE 7 (Tech Stack / Roadmap) ================
    (6, 0): "// 06  TECH STACK + ROADMAP",
    (6, 1): "BUILT WITH \u2014 NEXT 3\u20136 MONTHS",
    (6, 6): "Gemini 3 Pro & Flash",
    (6, 8): "Advanced reasoning, dynamic structuring, low-latency, hackathon core",
    (6, 11): "OBSERVABILITY",
    (6, 12): "Arize Phoenix Telemetry",
    (6, 14): "Trace spans, prompt parameters, cost tracking, dynamic override console",
    (6, 17): "DATA GROUNDING",
    (6, 18): "Elastic MCP & Tavily",
    (6, 20): "Retrieves macro indexes and news feeds for threat severity calculation",
    (6, 32): "Custom themed CSS, legibility selectors, responsive dashboard layout",
    (6, 38): "HQ Spoke-and-Hub Geopolitical maps, financial waterfall charts",
    (6, 44): "Live Elastic index syncing for geopolitical news",
    (6, 48): "Dynamic WACC calculations from localized interest rate feeds",
    (6, 52): "Trace evaluation dashboard inside Arize cloud UI",
    (6, 56): "PDF distress reports and export memos",
    (6, 60): "Interactive scenario debugger inside Streamlit",
    (6, 64): "Automatic baseline alerts when global chaos > 0.8",
    (6, 68): "Slack / Discord webhook alerts for portfolio assets",
    (6, 71): '"Doomsday Rapid Agent does not just guess systemic risk. It debates threats through an adversarial tribunal, grounds them via Elastic MCP, traces them via Arize, and models distress cascades in seconds."',
}

# Apply replacements
for (slide_idx, shape_idx), text in replacements.items():
    slide = prs.slides[slide_idx]
    shape = slide.shapes[shape_idx]
    set_shape_text(shape, text)

# Save
output_path = r'c:\Users\Moosa\Downloads\Doomsday_Rapid_Agent\Doomsday_Rapid_Agent.pptx'
prs.save(output_path)
print(f"Presentation saved to {output_path}")
