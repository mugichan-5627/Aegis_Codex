import pptx
import sys

prs = pptx.Presentation(r'c:\Users\Moosa\Downloads\Project_Doomsday.pptx')
out = []
for idx, slide in enumerate(prs.slides):
    out.append(f"\n================ SLIDE {idx+1} ================")
    for shape_idx, shape in enumerate(slide.shapes):
        if shape.has_text_frame and shape.text.strip():
            out.append(f"Shape {shape_idx} ({shape.name}):")
            for p_idx, p in enumerate(shape.text_frame.paragraphs):
                if p.text.strip():
                    out.append(f"  P{p_idx}: {p.text}")
        elif shape.has_table:
            out.append(f"Shape {shape_idx} ({shape.name}) is a Table:")
            table = shape.table
            for row_idx, row in enumerate(table.rows):
                row_txt = []
                for col_idx, cell in enumerate(row.cells):
                    row_txt.append(cell.text.strip())
                out.append(f"  Row {row_idx}: {row_txt}")

with open('doomsday_shape_details.txt', 'w', encoding='utf-8') as f:
    f.write('\n'.join(out))
print("Details dumped to doomsday_shape_details.txt")
